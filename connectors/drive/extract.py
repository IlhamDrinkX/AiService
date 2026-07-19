"""
Text extraction per Drive file type.

Google Docs are exported wholesale by the Drive API (files.export ->
text/plain) — that's already the complete document, no pagination/tabs
involved, nothing extra to fetch. Google Sheets and Google Slides are NOT
handled this way anymore (see below): Drive's export API only ever exposes
a subset of what's actually in those files, so both go through their own
API instead.

Everything else (PDF, docx, xlsx, pptx, rtf, txt, csv) is downloaded as raw
bytes and parsed locally. Images and other binary formats are left alone
here on purpose — that's the job of the separate photo-parsing pipeline
(vision model via Router AI), not this connector.

2026-07-19: extended for "collect maximum data" —
- Google Sheets: files.export(mimeType=text/csv) only ever returns the
  FIRST sheet/tab (a CSV can't represent multiple tabs) — multi-tab
  spreadsheets were silently losing every tab but the first. Now reads all
  tabs via the Sheets API (spreadsheets.get for the tab list, then
  spreadsheets.values.get per tab).
- Google Slides: files.export(mimeType=text/plain) only returns the
  visible on-slide text, never speaker notes. Now reads both via the
  Slides API.
- Added .pptx (python-pptx) and .rtf (striprtf) as extractable file types
  — previously silently skipped as "unsupported".
"""

import io
import logging

from googleapiclient.http import MediaIoBaseDownload

from retry import retry_call

log = logging.getLogger("drive_connector.extract")

# Google Docs only — Sheets/Slides are handled via their own APIs (see
# _sheets_text/_slides_text) because Drive's export endpoint can't expose
# their full content (multiple tabs / speaker notes respectively).
GOOGLE_NATIVE_EXPORT = {
    "application/vnd.google-apps.document": "text/plain",
}

GOOGLE_SHEETS_MIME = "application/vnd.google-apps.spreadsheet"
GOOGLE_SLIDES_MIME = "application/vnd.google-apps.presentation"

EXTRACTABLE_MIMES = set(GOOGLE_NATIVE_EXPORT) | {
    GOOGLE_SHEETS_MIME,
    GOOGLE_SLIDES_MIME,
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/rtf",
    "text/rtf",
    "text/plain",
    "text/csv",
}


# Smaller chunks = shorter-lived requests = less exposure to a flaky
# connection dropping mid-transfer, and a failed chunk only needs to redo
# a few MB instead of the whole file.
DOWNLOAD_CHUNK_SIZE = 5 * 1024 * 1024


def _download(request) -> bytes:
    buf = io.BytesIO()
    downloader = MediaIoBaseDownload(buf, request, chunksize=DOWNLOAD_CHUNK_SIZE)
    done = False
    while not done:
        # Capped retries/backoff on purpose: on a genuinely bad connection a
        # single stuck file shouldn't be able to block the whole run for
        # 10+ minutes — 3 attempts x ~60s timeout + short backoff, then give
        # up on this file and move on (it'll retry again on the next sync
        # since modified_time won't have been recorded for it).
        _, done = retry_call(downloader.next_chunk, what="media download", retries=3, base_delay=5)
    return buf.getvalue()


def _pdf_text(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _docx_text(data: bytes) -> str:
    import docx

    doc = docx.Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs)


def _xlsx_text(data: bytes) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines = []
    for ws in wb.worksheets:
        lines.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            lines.append("\t".join("" if c is None else str(c) for c in row))
    return "\n".join(lines)


def _pptx_text(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    lines = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    lines.append("\t".join(cell.text for cell in row.cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            lines.append(f"Notes: {slide.notes_slide.notes_text_frame.text}")
    return "\n".join(lines)


def _rtf_text(data: bytes) -> str:
    from striprtf.striprtf import rtf_to_text

    return rtf_to_text(data.decode("utf-8", errors="replace"))


def _sheets_text(sheets_service, file_id: str, name: str) -> str:
    """All tabs, not just the first — see module docstring for why this
    can't go through Drive's export API."""
    meta = retry_call(
        lambda: sheets_service.spreadsheets()
        .get(spreadsheetId=file_id, fields="sheets.properties.title")
        .execute(),
        what=f"sheets.get metadata ({name})",
    )
    titles = [s["properties"]["title"] for s in meta.get("sheets", [])]

    lines = []
    for title in titles:
        lines.append(f"# Sheet: {title}")
        # Sheet titles with spaces/special chars need single-quoting in an
        # A1-notation range; a bare title as the whole range means "all
        # used cells on that sheet".
        rng = f"'{title}'" if any(c in title for c in " !'\"") else title
        try:
            resp = retry_call(
                lambda rng=rng: sheets_service.spreadsheets()
                .values()
                .get(spreadsheetId=file_id, range=rng, valueRenderOption="FORMATTED_VALUE")
                .execute(),
                what=f"sheets.values.get ({name} / {title})",
            )
        except Exception as e:
            log.warning("  sheet '%s' in %s failed: %s", title, name, e)
            continue
        for row in resp.get("values", []):
            lines.append("\t".join(str(c) for c in row))
    return "\n".join(lines)


def _slides_text(slides_service, file_id: str, name: str) -> str:
    """On-slide text AND speaker notes — Drive's plain-text export only
    ever gives the former, see module docstring."""
    pres = retry_call(
        lambda: slides_service.presentations().get(presentationId=file_id).execute(),
        what=f"slides.get ({name})",
    )
    lines = []
    for i, slide in enumerate(pres.get("slides", []), start=1):
        lines.append(f"# Slide {i}")
        for el in slide.get("pageElements", []):
            shape = el.get("shape")
            if not shape:
                continue
            for te in shape.get("text", {}).get("textElements", []):
                run = te.get("textRun")
                if run and run.get("content"):
                    lines.append(run["content"].rstrip("\n"))
        notes_id = slide.get("slideProperties", {}).get("notesPage", {}).get("objectId")
        notes_page = slide.get("slideProperties", {}).get("notesPage")
        if notes_page:
            note_lines = []
            for el in notes_page.get("pageElements", []):
                shape = el.get("shape")
                if not shape or shape.get("placeholder", {}).get("type") != "BODY":
                    continue
                for te in shape.get("text", {}).get("textElements", []):
                    run = te.get("textRun")
                    if run and run.get("content"):
                        note_lines.append(run["content"].rstrip("\n"))
            note_text = "\n".join(t for t in note_lines if t.strip())
            if note_text.strip():
                lines.append(f"Notes: {note_text}")
    return "\n".join(lines)


def _cap(text: str, max_chars: int, name: str) -> str:
    """max_chars <= 0 means "no limit" — return the full text untouched.
    Documents (PDFs, Docs, etc.) longer than a nominal "50 pages" must be
    stored in full by default; TEXT_MAX_CHARS is an opt-in safety valve for
    someone who deliberately wants to cap storage/context size, not a
    default ceiling. When it does cut something, log it loudly — silent
    truncation is exactly the failure mode this is here to avoid."""
    if max_chars and max_chars > 0 and len(text) > max_chars:
        log.warning("Truncating %s: %d chars -> %d (TEXT_MAX_CHARS)", name, len(text), max_chars)
        return text[:max_chars]
    return text


def extract_text(services: dict, file_meta: dict, max_chars: int) -> str | None:
    """services: {"drive": ..., "sheets": ..., "slides": ...} — see auth.get_services()."""
    mime = file_meta["mimeType"]
    file_id = file_meta["id"]
    name = file_meta.get("name", file_id)
    drive = services["drive"]

    try:
        if mime == GOOGLE_SHEETS_MIME:
            return _cap(_sheets_text(services["sheets"], file_id, name), max_chars, name)

        if mime == GOOGLE_SLIDES_MIME:
            return _cap(_slides_text(services["slides"], file_id, name), max_chars, name)

        if mime in GOOGLE_NATIVE_EXPORT:
            request = drive.files().export_media(fileId=file_id, mimeType=GOOGLE_NATIVE_EXPORT[mime])
            return _cap(_download(request).decode("utf-8", errors="replace"), max_chars, name)

        if mime not in EXTRACTABLE_MIMES:
            return None  # unsupported (images, zips, binaries, etc.) — skip

        request = drive.files().get_media(fileId=file_id, supportsAllDrives=True)
        data = _download(request)

        if mime == "application/pdf":
            return _cap(_pdf_text(data), max_chars, name)
        if mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return _cap(_docx_text(data), max_chars, name)
        if mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return _cap(_xlsx_text(data), max_chars, name)
        if mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return _cap(_pptx_text(data), max_chars, name)
        if mime in ("application/rtf", "text/rtf"):
            return _cap(_rtf_text(data), max_chars, name)
        if mime in ("text/plain", "text/csv"):
            return _cap(data.decode("utf-8", errors="replace"), max_chars, name)

    except Exception as e:
        log.warning("Failed to extract text from %s (%s): %s", name, mime, e)
        return None

    return None
